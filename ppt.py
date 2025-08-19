from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 프레젠테이션 객체 생성
prs = Presentation()

# 슬라이드 레이아웃 선택 함수
def add_slide(title, content, layout=1):
    slide_layout = prs.slide_layouts[layout]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    # 제목 설정
    title_placeholder.text = title

    # 본문 내용 추가
    tf = content_placeholder.text_frame
    tf.clear()
    for line in content.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
    return slide

# 슬라이드 내용
slides_content = [
    ("📚🎬 프로젝트 소개", 
     "주제: 장르 기반 책 & 영화 추천 웹앱\n"
     "도구: Streamlit\n"
     "발표 시간: 4~5분"),
    
    ("❓ 문제 정의", 
     "• 책이나 영화를 고를 때 선택지가 너무 많아 고민됨\n"
     "• 장르별로 취향에 맞는 콘텐츠를 찾기 어려움\n"
     "• 최신 트렌드와 고전을 아우르는 추천 필요"),
    
    ("💡 솔루션", 
     "• Streamlit 기반 웹 앱 개발\n"
     "• 장르 선택 + 책/영화 선택 → 자동 추천\n"
     "• 포스터/표지 + 줄거리 제공"),
    
    ("⚙️ 개발 과정", 
     "1. Streamlit UI 제작 (selectbox, radio)\n"
     "2. 장르별 데이터베이스 구성 (책 + 영화)\n"
     "3. 이미지와 줄거리 추가\n"
     "4. 랜덤 추천 기능 구현"),
    
    ("✨ 주요 기능", 
     "• 4개 장르 제공 (로맨스, 판타지, 스릴러, SF)\n"
     "• 책과 영화 중 선택 가능\n"
     "• 최신 작품 + 고전 작품 균형 있게 제공\n"
     "• 3개 추천 결과 + 이미지/줄거리 출력"),
    
    ("📈 기대 효과", 
     "• 콘텐츠 선택 시간 절약\n"
     "• 독서/영화 문화 확산에 기여\n"
     "• 추천 시스템 확장 가능성 확보"),
    
    ("🚀 발전 방향", 
     "• AI 기반 개인화 추천 추가\n"
     "• 사용자 리뷰/평점 시스템 도입\n"
     "• 더 많은 장르와 작품 데이터 확장")
]

# 슬라이드 추가
for title, content in slides_content:
    add_slide(title, content)

# 파일 저장
prs.save("책_영화_추천앱_발표자료.pptx")
print("✅ PPTX 파일이 생성되었습니다: 책_영화_추천앱_발표자료.pptx")

